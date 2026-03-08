"""
Generate the data collection summary PowerPoint.

Covers all 7 data sources, row counts, API details,
merged output structure, and feature taxonomy.

Run from project root:
    python scripts/build_data_collection_slides.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1A, 0x2E, 0x4A)
GOLD  = RGBColor(0xC4, 0x9A, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xF4, 0xF4, 0xF4)
DGRAY = RGBColor(0x44, 0x44, 0x44)
MGRAY = RGBColor(0x88, 0x88, 0x88)
TEAL  = RGBColor(0x00, 0x7A, 0x87)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
AMBER = RGBColor(0xD4, 0x7A, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

DOCS_DIR = Path(__file__).parent.parent / "docs"


# ── Primitives ────────────────────────────────────────────────────────────────

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, line=None, line_w=Pt(0.75)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def txb(slide, l, t, w, h, text, size, color=WHITE, bold=False,
        italic=False, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


def title_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, SLIDE_W, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.09), SLIDE_W, Inches(0.06), GOLD)
    txb(slide, Inches(0.45), Inches(0.17), Inches(12.4), Inches(0.65),
        title, 28, WHITE, bold=True)
    if subtitle:
        txb(slide, Inches(0.45), Inches(0.80), Inches(12.4), Inches(0.30),
            subtitle, 13, MGRAY, italic=True)


def footer(slide, text="CMU MSBA Capstone · Stablecoin Depeg Prediction · March 2026"):
    rect(slide, 0, SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32), NAVY)
    txb(slide, Inches(0.3), SLIDE_H - Inches(0.30), Inches(12.7), Inches(0.28),
        text, 9, MGRAY)


def section_divider(title, subtitle=""):
    sl = prs.slides.add_slide(BLANK)
    bg(sl, NAVY)
    rect(sl, 0, Inches(3.0), SLIDE_W, Inches(0.08), GOLD)
    txb(sl, Inches(1.2), Inches(1.8), Inches(10.9), Inches(1.2),
        title, 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        txb(sl, Inches(1.2), Inches(3.2), Inches(10.9), Inches(0.5),
            subtitle, 18, GOLD, align=PP_ALIGN.CENTER)
    footer(sl)


def status_chip(slide, l, t, label, color, w=Inches(1.3)):
    rect(slide, l, t, w, Inches(0.28), color)
    txb(slide, l, t + Inches(0.03), w, Inches(0.25),
        label, 10, WHITE, bold=True, align=PP_ALIGN.CENTER)


def table_header(slide, y, col_x, col_w, headers, row_h=Inches(0.44)):
    rect(slide, col_x[0], y, sum(col_w) + (col_x[-1] - col_x[0] - sum(col_w[:-1])), row_h, NAVY)
    for cx, cw, h in zip(col_x, col_w, headers):
        txb(slide, cx + Inches(0.06), y + Inches(0.11), cw, Inches(0.26), h, 11.5, WHITE, bold=True)
    return y + row_h


def table_row(slide, y, col_x, col_w, vals, colors, bolds, bg_col, row_h=Inches(0.44)):
    total_w = col_x[-1] + col_w[-1] - col_x[0]
    rect(slide, col_x[0], y, total_w, row_h, bg_col, MGRAY)
    for cx, cw, val, fc, fb in zip(col_x, col_w, vals, colors, bolds):
        txb(slide, cx + Inches(0.06), y + Inches(0.11), cw, Inches(0.26), val, 11, fc, bold=fb)
    return y + row_h


# ── Slide 1 — Title ───────────────────────────────────────────────────────────

def slide_title():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, NAVY)
    rect(sl, 0, Inches(3.05), SLIDE_W, Inches(0.08), GOLD)
    txb(sl, Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.1),
        "Data Collection Summary", 46, WHITE, bold=True, align=PP_ALIGN.CENTER)
    txb(sl, Inches(0.8), Inches(2.55), Inches(11.7), Inches(0.55),
        "Stablecoin Depeg Prediction at 5-Minute Resolution", 24, GOLD, align=PP_ALIGN.CENTER)
    txb(sl, Inches(0.8), Inches(3.25), Inches(11.7), Inches(0.4),
        "CMU MSBA Capstone · March 2026", 16, MGRAY, align=PP_ALIGN.CENTER)
    txb(sl, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.35),
        "Robert Springett", 18, WHITE, align=PP_ALIGN.CENTER)
    footer(sl)


# ── Slide 2 — Overview ────────────────────────────────────────────────────────

def slide_overview():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, LGRAY)
    title_bar(sl, "Data Collection Overview",
              "All sources collected at 5-minute resolution or forward-filled from daily")
    footer(sl)

    sources = [
        (GREEN,  "1",  "Binance OHLCV",             "Free · 5m native",   "BTC/ETH market context for all coins"),
        (GREEN,  "2",  "CoinAPI VWAP Index",         "Paid · 5m native",   "Primary peg price signal · cross-exchange VWAP"),
        (GREEN,  "3",  "ETH On-Chain Mint/Burn",     "Free · event-level", "Supply creation/destruction + USDT treasury flows"),
        (GREEN,  "4",  "USDT TRON Treasury",         "Free · event-level", "Redemption pressure — >50% of USDT supply lives on TRON"),
        (GREEN,  "5",  "Curve Pool Swaps",           "Free · event-level", "DEX stress indicator — leading depeg signal"),
        (GREEN,  "6",  "FRED Macro",                 "Free · daily",       "DXY, VIX, 10Y yield, Fed Funds → forward-filled"),
        (GREEN,  "7",  "Market Daily",               "Free · daily",       "Fear & Greed index → forward-filled"),
    ]

    col_x = [Inches(0.30), Inches(0.75), Inches(2.40), Inches(5.30), Inches(7.10)]
    col_w = [Inches(0.38), Inches(1.55), Inches(2.80), Inches(1.70), Inches(5.90)]
    row_h = Inches(0.62)
    y = Inches(1.35)

    for scol, num, name, freq, detail in sources:
        rect(sl, col_x[0], y, Inches(12.75), row_h, WHITE, MGRAY)
        rect(sl, col_x[0], y, Inches(0.08), row_h, scol)
        txb(sl, Inches(0.50), y + Inches(0.16), Inches(0.30), Inches(0.28), num, 13, GOLD, bold=True)
        txb(sl, col_x[1] + Inches(0.06), y + Inches(0.05), col_w[1], Inches(0.28), name, 13, NAVY, bold=True)
        txb(sl, col_x[1] + Inches(0.06), y + Inches(0.33), col_w[1], Inches(0.24), freq, 10, TEAL, italic=True)
        txb(sl, col_x[2] + Inches(0.06), y + Inches(0.05), col_w[2], Inches(0.28), "Complete", 11, GREEN, bold=True)
        txb(sl, col_x[3] + Inches(0.06), y + Inches(0.16), col_w[3] + col_w[4], Inches(0.28), detail, 11.5, DGRAY)
        y += row_h + Inches(0.06)


# ── Slide 3 — Stablecoins Covered ────────────────────────────────────────────

def slide_coins():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, LGRAY)
    title_bar(sl, "Stablecoins Covered",
              "Failed and discontinued coins included — they provide real depeg training examples")
    footer(sl)

    col_x = [Inches(0.30), Inches(1.35), Inches(3.35), Inches(5.80), Inches(7.75), Inches(9.85), Inches(11.65)]
    col_w = [Inches(0.95), Inches(1.90), Inches(2.35), Inches(1.85), Inches(2.00), Inches(1.70), Inches(1.50)]
    headers = ["Coin", "Full Name", "Type", "Status", "Coverage Start", "Coverage End", "5m Rows"]
    row_h = Inches(0.50)

    rows = [
        ("USDT",  "Tether",        "Fiat-backed",           "Active",        "Aug 2017", "Feb 2026", "897,984",  GREEN),
        ("USDC",  "USD Coin",      "Fiat-backed",           "Active",        "Oct 2018", "Feb 2026", "775,584",  GREEN),
        ("DAI",   "MakerDAO Dai",  "Crypto-collat.",        "Active",        "Apr 2018", "Feb 2026", "828,963",  GREEN),
        ("BUSD",  "Binance USD",   "Fiat-backed",           "Discontinued",  "Sep 2019", "Mar 2023", "373,609",  AMBER),
        ("UST",   "TerraUSD",      "Algorithmic",           "Failed",        "Nov 2020", "May 2022", "159,433",  RGBColor(0xC0,0x3B,0x3B)),
        ("USDe",  "Ethena USDe",   "Synthetic",             "Active",        "Apr 2024", "Feb 2026", "201,024",  GREEN),
        ("RLUSD", "Ripple USD",    "Fiat-backed",           "Active",        "Feb 2025", "Feb 2026", "105,696",  GREEN),
    ]

    y = Inches(1.32)
    rect(sl, col_x[0], y, Inches(12.85), row_h, NAVY)
    for cx, cw, h in zip(col_x, col_w, headers):
        txb(sl, cx + Inches(0.05), y + Inches(0.12), cw, Inches(0.28), h, 11.5, WHITE, bold=True)
    y += row_h

    for i, (ticker, name, ctype, status, start, end, rows_n, scol) in enumerate(rows):
        bg_c = WHITE if i % 2 == 0 else RGBColor(0xEB, 0xF2, 0xFA)
        rect(sl, col_x[0], y, Inches(12.85), row_h, bg_c, MGRAY)
        txb(sl, col_x[0]+Inches(0.05), y+Inches(0.12), col_w[0], Inches(0.28), ticker, 12, NAVY, bold=True)
        txb(sl, col_x[1]+Inches(0.05), y+Inches(0.12), col_w[1], Inches(0.28), name,   11, DGRAY)
        txb(sl, col_x[2]+Inches(0.05), y+Inches(0.12), col_w[2], Inches(0.28), ctype,  10, MGRAY, italic=True)
        status_chip(sl, col_x[3]+Inches(0.05), y+Inches(0.12), status, scol, w=Inches(1.55))
        txb(sl, col_x[4]+Inches(0.05), y+Inches(0.12), col_w[4], Inches(0.28), start,  11, DGRAY)
        txb(sl, col_x[5]+Inches(0.05), y+Inches(0.12), col_w[5], Inches(0.28), end,    11, DGRAY)
        txb(sl, col_x[6]+Inches(0.05), y+Inches(0.12), col_w[6], Inches(0.28), rows_n, 11, DGRAY)
        y += row_h

    txb(sl, Inches(0.35), Inches(5.75), Inches(12.6), Inches(0.28),
        "Depeg threshold: |coinapi_close − 1.00| > 0.005  for ≥ 3 consecutive 5-min bars (15 min sustained)",
        11.5, TEAL, italic=True)


# ── Slide 4 — CoinAPI VWAP ───────────────────────────────────────────────────

def slide_coinapi():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, LGRAY)
    title_bar(sl, "Source 2 — CoinAPI VWAP Index",
              "IDX_REFRATE_VWAP — cross-exchange volume-weighted average price · Paid API")
    footer(sl)

    # Left: key facts
    rect(sl, Inches(0.30), Inches(1.35), Inches(6.0), Inches(5.70), WHITE, MGRAY)
    txb(sl, Inches(0.50), Inches(1.45), Inches(5.6), Inches(0.28), "Why CoinAPI VWAP?", 13, NAVY, bold=True)
    bullets = [
        "Primary peg measurement signal — aggregates across all major venues",
        "Eliminates single-exchange noise and manipulation risk",
        "Native 5-minute resolution — no aggregation needed",
        "Covers fiat pairs (USD) not available on crypto-only exchanges",
        "Used as coinapi_open/high/low/close in all processed files",
    ]
    by = Inches(1.80)
    for b in bullets:
        rect(sl, Inches(0.52), by + Inches(0.08), Inches(0.06), Inches(0.06), TEAL)
        txb(sl, Inches(0.70), by, Inches(5.4), Inches(0.30), b, 11, DGRAY)
        by += Inches(0.36)

    txb(sl, Inches(0.50), Inches(3.70), Inches(5.6), Inches(0.28), "API Details", 13, NAVY, bold=True)
    api_details = [
        ("Endpoint",   "rest-api.indexes.coinapi.io/v1/indexes/{id}/timeseries"),
        ("Params",     "period_id=5MIN, time_start, time_end, limit=100000"),
        ("Auth",       "COINAPI_KEY required (paid subscription)"),
        ("Rate limit", "~10 req/s on standard plan; checkpoint/resume built in"),
    ]
    dy = Inches(4.02)
    for label, val in api_details:
        txb(sl, Inches(0.52), dy, Inches(1.10), Inches(0.26), label, 10.5, NAVY, bold=True)
        txb(sl, Inches(1.68), dy, Inches(4.4), Inches(0.26), val, 10.5, DGRAY)
        dy += Inches(0.30)

    # Right: per-coin table
    col_x = [Inches(6.55), Inches(7.85), Inches(10.05)]
    col_w = [Inches(1.20), Inches(2.10), Inches(3.10)]
    headers = ["Coin", "5m Rows", "Date Range"]
    row_h = Inches(0.44)

    tbl_rows = [
        ("USDT",  "897,984",  "Aug 2017 → Feb 2026"),
        ("USDC",  "775,584",  "Oct 2018 → Feb 2026"),
        ("DAI",   "828,963",  "Apr 2018 → Feb 2026"),
        ("BUSD",  "373,609",  "Sep 2019 → Mar 2023"),
        ("UST",   "159,433",  "Nov 2020 → May 2022"),
        ("USDe",  "201,024",  "Apr 2024 → Feb 2026"),
        ("RLUSD", "105,696",  "Feb 2025 → Feb 2026"),
        ("TOTAL", "3,342,293","Aug 2017 → Feb 2026"),
    ]

    y = Inches(1.35)
    rect(sl, col_x[0], y, Inches(6.50), row_h, NAVY)
    for cx, cw, h in zip(col_x, col_w, headers):
        txb(sl, cx + Inches(0.06), y + Inches(0.11), cw, Inches(0.26), h, 11.5, WHITE, bold=True)
    y += row_h

    for i, (coin, recs, dr) in enumerate(tbl_rows):
        is_total = coin == "TOTAL"
        bg_c = NAVY if is_total else (WHITE if i % 2 == 0 else RGBColor(0xEB,0xF2,0xFA))
        fc = WHITE if is_total else DGRAY
        rect(sl, col_x[0], y, Inches(6.50), row_h, bg_c, MGRAY)
        txb(sl, col_x[0]+Inches(0.06), y+Inches(0.11), col_w[0], Inches(0.26),
            coin, 12, WHITE if is_total else NAVY, bold=True)
        txb(sl, col_x[1]+Inches(0.06), y+Inches(0.11), col_w[1], Inches(0.26), recs, 11, fc, bold=is_total)
        txb(sl, col_x[2]+Inches(0.06), y+Inches(0.11), col_w[2], Inches(0.26), dr,   11, fc)
        y += row_h


# ── Slide 5 — ETH On-Chain ───────────────────────────────────────────────────

def slide_onchain():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, LGRAY)
    title_bar(sl, "Source 3 — ETH On-Chain Mint/Burn + USDT Treasury",
              "Etherscan V2 Logs API · event-level → aggregated to 5-min bins")
    footer(sl)

    # Mint/burn table
    txb(sl, Inches(0.35), Inches(1.30), Inches(12.6), Inches(0.28),
        "Mint / Burn Events by Coin", 13, NAVY, bold=True)

    col_x = [Inches(0.30), Inches(1.35), Inches(3.10), Inches(5.50), Inches(7.80), Inches(10.15)]
    col_w = [Inches(0.95), Inches(1.65), Inches(2.30), Inches(2.20), Inches(2.25), Inches(2.90)]
    headers = ["Coin", "Raw Events", "5m Rows", "Date Range", "Mint Event", "Burn Event"]
    row_h = Inches(0.42)

    rows = [
        ("USDT",  "11,215",    "8,703",   "Nov 2017 → Feb 2026", "Issue(uint256)",              "DestroyedBlackFunds (OFAC)"),
        ("USDC",  "2,209,692", "390,104", "Sep 2018 → Mar 2026", "Mint(address,address,uint256)","Burn(address,uint256)"),
        ("DAI",   "1,400,464", "425,580", "Nov 2019 → Mar 2026", "Transfer(from=0x0)",           "Transfer(to=0x0)"),
        ("BUSD",  "3,767",     "3,765",   "Sep 2019 → Jan 2026", "SupplyIncreased(address,uint)","SupplyDecreased(address,uint)"),
        ("USDe",  "38,001",    "24,589",  "Nov 2023 → Mar 2026", "Transfer(from=0x0)",           "Transfer(to=0x0)"),
        ("RLUSD", "288",       "282",     "Aug 2024 → Feb 2026", "Transfer(from=0x0)",           "Transfer(to=0x0)"),
    ]

    y = Inches(1.60)
    rect(sl, col_x[0], y, Inches(13.0), row_h, NAVY)
    for cx, cw, h in zip(col_x, col_w, headers):
        txb(sl, cx+Inches(0.04), y+Inches(0.10), cw, Inches(0.26), h, 10.5, WHITE, bold=True)
    y += row_h

    for i, (coin, evts, bins, dr, mint, burn) in enumerate(rows):
        bg_c = WHITE if i % 2 == 0 else RGBColor(0xEB, 0xF2, 0xFA)
        rect(sl, col_x[0], y, Inches(13.0), row_h, bg_c, MGRAY)
        for cx, cw, val, fc, fb in zip(col_x, col_w,
                [coin, evts, bins, dr, mint, burn],
                [NAVY, DGRAY, DGRAY, DGRAY, MGRAY, MGRAY],
                [True, False, False, False, False, False]):
            txb(sl, cx+Inches(0.04), y+Inches(0.10), cw, Inches(0.26), val, 10, fc, bold=fb)
        y += row_h

    # USDT treasury note
    rect(sl, Inches(0.30), Inches(4.80), Inches(12.75), Inches(1.15), NAVY)
    rect(sl, Inches(0.30), Inches(4.80), Inches(0.08), Inches(1.15), GOLD)
    txb(sl, Inches(0.52), Inches(4.87), Inches(12.2), Inches(0.26),
        "USDT ETH Treasury Flows (included in usdt_eth_5m.parquet)", 12, GOLD, bold=True)
    txb(sl, Inches(0.52), Inches(5.16), Inches(12.2), Inches(0.68),
        "Institutions return USDT to 0x5754284f… (Tether Treasury) for off-chain USD redemption. "
        "These inflows (treasury_inflow_*) are a direct redemption-pressure signal not visible in "
        "mint/burn data. Tether never calls redeem() on Ethereum — no on-chain burns for real redemptions.",
        11, WHITE)

    txb(sl, Inches(0.35), Inches(6.05), Inches(12.6), Inches(0.26),
        "5m bin columns: mint_count, mint_volume_usd, burn_count, burn_volume_usd, net_flow_usd  "
        "|  USDT also: treasury_inflow_count/volume, treasury_outflow_count/volume, treasury_net_flow_usd",
        10, TEAL, italic=True)


# ── Slide 6 — TRON + Curve ───────────────────────────────────────────────────

def slide_tron_curve():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, LGRAY)
    title_bar(sl, "Sources 4 & 5 — USDT TRON Treasury + Curve Pool Swaps",
              "Redemption pressure on TRON · DEX stress leading indicator")
    footer(sl)

    # TRON box (left)
    rect(sl, Inches(0.30), Inches(1.35), Inches(6.1), Inches(5.70), WHITE, MGRAY)
    rect(sl, Inches(0.30), Inches(1.35), Inches(0.08), Inches(5.70), TEAL)
    txb(sl, Inches(0.52), Inches(1.42), Inches(5.7), Inches(0.28),
        "USDT TRON Treasury Flows", 14, NAVY, bold=True)
    txb(sl, Inches(0.52), Inches(1.73), Inches(5.7), Inches(0.55),
        "TRON hosts >50% of USDT supply — where Tether actively manages mint/burn. "
        "Treasury inflow = institution returning USDT for USD redemption.",
        11, DGRAY)

    tron_facts = [
        ("API",       "TronGrid REST · TRONGRID_API_KEY optional"),
        ("Wallets",   "3 treasury wallets tracked (2019–present)"),
        ("Transfers", "12,663 total · inter-treasury excluded"),
        ("5m bins",   "9,610 rows · Apr 2019 → Feb 2026"),
        ("Columns",   "tron_treasury_inflow/outflow count + volume + net"),
    ]
    dy = Inches(2.35)
    for label, val in tron_facts:
        txb(sl, Inches(0.52), dy, Inches(1.1), Inches(0.26), label, 10.5, NAVY, bold=True)
        txb(sl, Inches(1.70), dy, Inches(4.5), Inches(0.26), val, 10.5, DGRAY)
        dy += Inches(0.32)

    txb(sl, Inches(0.52), Inches(3.90), Inches(5.7), Inches(0.28),
        "Hypothesis", 12, NAVY, bold=True)
    txb(sl, Inches(0.52), Inches(4.20), Inches(5.7), Inches(0.80),
        "Sophisticated institutions redeem at par directly with Tether before market depegs. "
        "TRON treasury inflow spikes may precede open-market peg breaks by minutes to hours.",
        11, DGRAY)

    # Curve box (right)
    rect(sl, Inches(6.70), Inches(1.35), Inches(6.35), Inches(5.70), WHITE, MGRAY)
    rect(sl, Inches(6.70), Inches(1.35), Inches(0.08), Inches(5.70), GOLD)
    txb(sl, Inches(6.92), Inches(1.42), Inches(5.9), Inches(0.28),
        "Curve Pool Swaps", 14, NAVY, bold=True)
    txb(sl, Inches(6.92), Inches(1.73), Inches(5.9), Inches(0.40),
        "TokenExchange events on 3 pools aggregated to 5-min sold/bought/net per token.",
        11, DGRAY)

    curve_pools = [
        ("3pool",      "DAI/USDC/USDT", "Sep 2020", "538,925", "277,828"),
        ("usde_usdc",  "USDe/USDC",     "Nov 2023", "106,449",  "61,095"),
        ("rlusd_usdc", "RLUSD/USDC",    "Dec 2024",  "10,719",   "8,835"),
    ]
    col_x2 = [Inches(6.92), Inches(8.30), Inches(9.55), Inches(10.60), Inches(11.75)]
    col_w2 = [Inches(1.28), Inches(1.15), Inches(0.95), Inches(1.05), Inches(1.10)]
    heads2 = ["Pool", "Tokens", "Since", "Events", "5m bins"]
    y = Inches(2.24)
    rect(sl, col_x2[0], y, Inches(6.15), Inches(0.36), NAVY)
    for cx, cw, h in zip(col_x2, col_w2, heads2):
        txb(sl, cx+Inches(0.04), y+Inches(0.06), cw, Inches(0.24), h, 10, WHITE, bold=True)
    y += Inches(0.36)
    for i, (pool, tokens, since, evts, bins) in enumerate(curve_pools):
        bg_c = WHITE if i % 2 == 0 else RGBColor(0xEB,0xF2,0xFA)
        rect(sl, col_x2[0], y, Inches(6.15), Inches(0.40), bg_c, MGRAY)
        for cx, cw, val, fc in zip(col_x2, col_w2,
                [pool, tokens, since, evts, bins],
                [NAVY, DGRAY, DGRAY, DGRAY, DGRAY]):
            txb(sl, cx+Inches(0.04), y+Inches(0.08), cw, Inches(0.26), val, 10.5, fc)
        y += Inches(0.40)

    txb(sl, Inches(6.92), Inches(3.55), Inches(5.9), Inches(0.28),
        "Key Signal", 12, NAVY, bold=True)
    txb(sl, Inches(6.92), Inches(3.85), Inches(5.9), Inches(0.90),
        "usdt_net_sell_volume_usd in 3pool > 0 means USDT being sold for "
        "DAI/USDC — sophisticated traders exiting on-chain before CEX prices move. "
        "Historically precedes open-market depeg by minutes.",
        11, DGRAY)

    txb(sl, Inches(6.92), Inches(4.85), Inches(5.9), Inches(0.28),
        "Columns per token (each pool):", 11, NAVY, bold=True)
    txb(sl, Inches(6.92), Inches(5.15), Inches(5.9), Inches(0.60),
        "{token}_sold_count, {token}_sold_volume_usd\n"
        "{token}_bought_count, {token}_bought_volume_usd\n"
        "{token}_net_sell_volume_usd",
        10.5, DGRAY)


# ── Slide 7 — Macro + Market ──────────────────────────────────────────────────

def slide_macro():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, LGRAY)
    title_bar(sl, "Sources 6 & 7 — FRED Macro + Market Daily",
              "Daily sources forward-filled to 5-minute resolution in merge pipeline")
    footer(sl)

    for panel_x, panel_w, title, color, rows_data in [
        (Inches(0.30), Inches(6.1), "FRED Macro Data", TEAL, [
            ("API",        "api.stlouisfed.org/fred/series/observations"),
            ("Auth",       "FRED_API_KEY required (free)"),
            ("Records",    "2,892 daily rows · Jan 2015 → Feb 2026"),
            ("Series",     "DXY, VIX, T10Y, FEDFUNDS"),
            ("In merge",   "Reindexed to 5m UTC, forward-filled"),
        ]),
        (Inches(6.90), Inches(6.15), "Market Daily", GOLD, [
            ("Sources",    "CoinGecko (BTC/ETH prices) + Alternative.me (Fear & Greed)"),
            ("Auth",       "COINGECKO_API_KEY optional · alternative.me free"),
            ("Records",    "2,949 daily rows · Feb 2018 → Feb 2026"),
            ("Series",     "fear_greed index (0–100)"),
            ("In merge",   "Reindexed to 5m UTC, forward-filled"),
        ]),
    ]:
        rect(sl, panel_x, Inches(1.35), panel_w, Inches(4.70), WHITE, MGRAY)
        rect(sl, panel_x, Inches(1.35), Inches(0.08), Inches(4.70), color)
        txb(sl, panel_x + Inches(0.22), Inches(1.42), panel_w - Inches(0.3), Inches(0.28),
            title, 14, NAVY, bold=True)
        dy = Inches(1.82)
        for label, val in rows_data:
            txb(sl, panel_x + Inches(0.22), dy, Inches(1.10), Inches(0.26), label, 10.5, NAVY, bold=True)
            txb(sl, panel_x + Inches(1.40), dy, panel_w - Inches(1.55), Inches(0.26), val, 10.5, DGRAY)
            dy += Inches(0.38)

    # Note on BTC/ETH
    rect(sl, Inches(0.30), Inches(6.20), Inches(12.75), Inches(0.86), NAVY)
    rect(sl, Inches(0.30), Inches(6.20), Inches(0.08), Inches(0.86), GOLD)
    txb(sl, Inches(0.52), Inches(6.27), Inches(12.2), Inches(0.26),
        "BTC/ETH 5-Minute Context (Binance BTCUSDT/ETHUSDT)", 12, GOLD, bold=True)
    txb(sl, Inches(0.52), Inches(6.55), Inches(12.2), Inches(0.44),
        "binance_btc_close and binance_eth_close joined at 5-min resolution to all 7 coins. "
        "Replaces daily forward-filled BTC/ETH price with intraday crypto market regime signal — "
        "a leading indicator for risk-off events that precede stablecoin stress.",
        11, WHITE)


# ── Slide 8 — Merged Output ───────────────────────────────────────────────────

def slide_merged():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, LGRAY)
    title_bar(sl, "Merged Output — data/processed/{coin}_5m.parquet",
              "Per-coin files with all sources joined · depeg labels applied · cutoff 2026-02-28")
    footer(sl)

    col_x = [Inches(0.30), Inches(1.35), Inches(2.75), Inches(4.30), Inches(5.65),
             Inches(6.80), Inches(8.00), Inches(9.30), Inches(10.65)]
    col_w = [Inches(0.95), Inches(1.30), Inches(1.45), Inches(1.25), Inches(1.05),
             Inches(1.10), Inches(1.20), Inches(1.25), Inches(2.40)]
    headers = ["Coin", "5m Rows", "Columns", "CoinAPI", "On-chain",
               "TRON", "Curve", "BTC/ETH", "Depeg rate (3-bar)"]
    row_h = Inches(0.46)

    rows = [
        ("USDT",  "897,984", "60", "✓",  "✓ ETH+treas.", "✓", "3pool",       "✓", "9.09%"),
        ("USDC",  "775,584", "39", "✓",  "✓ mint/burn",  "—", "3pool",       "✓", "5.59%"),
        ("DAI",   "828,963", "39", "✓",  "✓ mint/burn",  "—", "3pool",       "✓", "19.30%"),
        ("BUSD",  "373,609", "24", "✓",  "✓ mint/burn",  "—", "—",           "✓", "0.29%"),
        ("UST",   "159,433", "18", "✓",  "—",            "—", "—",           "✓", "8.58%"),
        ("USDe",  "201,024", "34", "✓",  "✓ mint/burn",  "—", "usde_usdc",   "✓", "0.05%"),
        ("RLUSD", "105,696", "34", "✓",  "✓ mint/burn",  "—", "rlusd_usdc",  "✓", "0.00%"),
    ]

    y = Inches(1.32)
    rect(sl, col_x[0], y, Inches(12.85), row_h, NAVY)
    for cx, cw, h in zip(col_x, col_w, headers):
        txb(sl, cx+Inches(0.04), y+Inches(0.11), cw, Inches(0.26), h, 10.5, WHITE, bold=True)
    y += row_h

    for i, (coin, rows_n, cols, ca, oc, tron, curve, btceth, depeg) in enumerate(rows):
        bg_c = WHITE if i % 2 == 0 else RGBColor(0xEB, 0xF2, 0xFA)
        rect(sl, col_x[0], y, Inches(12.85), row_h, bg_c, MGRAY)
        dval = float(depeg.replace("%",""))
        dcol = AMBER if dval > 15 else (GREEN if dval < 1 else DGRAY)
        for cx, cw, val, fc, fb in zip(col_x, col_w,
                [coin, rows_n, cols, ca, oc, tron, curve, btceth, depeg],
                [NAVY, DGRAY, DGRAY, GREEN if ca=="✓" else MGRAY,
                 GREEN if "✓" in oc else MGRAY,
                 GREEN if tron=="✓" else MGRAY,
                 TEAL if curve != "—" else MGRAY,
                 GREEN if btceth=="✓" else MGRAY, dcol],
                [True,False,False,False,False,False,False,False,True]):
            txb(sl, cx+Inches(0.04), y+Inches(0.11), cw, Inches(0.26), val, 11, fc, bold=fb)
        y += row_h

    # Label definition footer
    rect(sl, Inches(0.30), Inches(5.85), Inches(12.75), Inches(1.22), NAVY)
    rect(sl, Inches(0.30), Inches(5.85), Inches(0.08), Inches(1.22), GOLD)
    txb(sl, Inches(0.52), Inches(5.92), Inches(12.2), Inches(0.26),
        "Depeg Label Definition + Forward-Looking Prediction Targets", 12, GOLD, bold=True)
    txb(sl, Inches(0.52), Inches(6.20), Inches(12.2), Inches(0.76),
        "depeg = 1  if  |coinapi_close − 1.00| > 0.005  for ≥ 3 consecutive bars  (15 min sustained)\n"
        "Forward labels: depeg_next_5min (1 bar)  ·  depeg_next_30min (6 bars)  ·  "
        "depeg_next_1h (12 bars)  ·  depeg_next_4h (48 bars)  |  price_dev = coinapi_close − 1.00",
        11, WHITE)


# ── Slide 9 — Feature Taxonomy ────────────────────────────────────────────────

def slide_features():
    sl = prs.slides.add_slide(BLANK)
    bg(sl, LGRAY)
    title_bar(sl, "Feature Taxonomy",
              "Raw columns available before feature engineering · source-prefixed naming")
    footer(sl)

    categories = [
        (TEAL,  "Peg Price",
         "coinapi_open/high/low/close, coinapi_tick_count",
         "CoinAPI VWAP cross-exchange reference rate — primary depeg signal"),
        (NAVY,  "Binance Market Pairs (USDT only)",
         "binance_usdcusdt_open/high/low/close/volume + buy_ratio, spread_proxy",
         "USDC/USDT spot pair on Binance — secondary price signal for USDT"),
        (GREEN, "On-Chain Supply",
         "mint_count, mint_volume_usd, burn_count, burn_volume_usd, net_flow_usd",
         "Institutional supply creation/destruction (USDC/DAI/USDe/RLUSD most informative)"),
        (GREEN, "Redemption Pressure (USDT)",
         "treasury_inflow/outflow count+volume+net  ·  tron_treasury_inflow/outflow count+volume+net",
         "USDT-specific: institutions returning USDT to Tether for USD redemption on ETH and TRON"),
        (AMBER, "DEX Stress",
         "curve_{pool}_{token}_sold/bought/net_sell_volume_usd  (per pool, per token)",
         "Curve pool imbalance — informed traders fleeing a coin on-chain before CEX depeg"),
        (TEAL,  "BTC/ETH Market Context",
         "binance_btc_close, binance_eth_close",
         "5-min crypto market regime — proxy for risk-off events that trigger stablecoin stress"),
        (DGRAY, "Macro Regime",
         "dxy, vix, t10y, fedfunds",
         "FRED daily macro indicators (forward-filled) — dollar strength, market stress, rates"),
        (DGRAY, "Sentiment",
         "fear_greed",
         "Alternative.me Fear & Greed index (forward-filled) — retail crypto sentiment"),
        (GOLD,  "Unified Flow Signal",
         "total_net_flow_usd",
         "Cross-coin comparable: USDT = ETH+TRON treasury net flow; others = mint/burn net flow"),
    ]

    y = Inches(1.32)
    row_h = Inches(0.58)
    for scol, cat, cols, note in categories:
        rect(sl, Inches(0.30), y, Inches(12.75), row_h, WHITE, MGRAY)
        rect(sl, Inches(0.30), y, Inches(0.08), row_h, scol)
        txb(sl, Inches(0.52), y + Inches(0.04), Inches(2.5), Inches(0.26), cat, 12, NAVY, bold=True)
        txb(sl, Inches(3.10), y + Inches(0.04), Inches(9.8), Inches(0.26), cols, 10.5, TEAL)
        txb(sl, Inches(0.52), y + Inches(0.30), Inches(12.3), Inches(0.24), note, 10.5, DGRAY)
        y += row_h + Inches(0.03)


# ── Build ─────────────────────────────────────────────────────────────────────

slide_title()
slide_overview()
slide_coins()
slide_coinapi()
slide_onchain()
slide_tron_curve()
slide_macro()
slide_merged()
slide_features()

# ── Save ──────────────────────────────────────────────────────────────────────

out = DOCS_DIR / "data_collection_summary.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
