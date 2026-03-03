"""
Build data_collection_summary.pptx
CMU MSBA Capstone — Stablecoin Depeg Prediction
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pathlib import Path

# ---------------------------------------------------------------------------
# Colours
# ---------------------------------------------------------------------------
CMU_RED    = RGBColor(0xC4, 0x12, 0x30)
DARK_GREY  = RGBColor(0x2D, 0x2D, 0x2D)
MID_GREY   = RGBColor(0x60, 0x60, 0x60)
LIGHT_GREY = RGBColor(0xF4, 0xF4, 0xF4)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)
AMBER      = RGBColor(0xE6, 0x5C, 0x00)
BLUE       = RGBColor(0x15, 0x65, 0xC0)

# ---------------------------------------------------------------------------
# Slide dimensions  (widescreen 13.33" × 7.5")
# ---------------------------------------------------------------------------
W = Inches(13.33)
H = Inches(7.50)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle=None):
    """Red top bar with white title text."""
    add_rect(slide, 0, 0, W, Inches(1.1), fill=CMU_RED)
    add_text(slide, title, Inches(0.35), Inches(0.12), Inches(11), Inches(0.65),
             size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, Inches(0.35), Inches(0.68), Inches(12), Inches(0.38),
                 size=13, bold=False, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.LEFT)
    # thin red bottom line
    add_rect(slide, 0, Inches(1.1), W, Pt(3), fill=RGBColor(0x8B, 0x0D, 0x22))


def footer(slide, page_num, total):
    add_text(slide, f"CMU MSBA Capstone  |  Stablecoin Depeg Prediction  |  {page_num}/{total}",
             Inches(0.3), Inches(7.15), Inches(12.5), Inches(0.3),
             size=9, color=MID_GREY, align=PP_ALIGN.CENTER)
    add_rect(slide, 0, Inches(7.1), W, Pt(1.5), fill=CMU_RED)


def table(slide, headers, rows, l, t, w, col_widths=None,
          header_fill=CMU_RED, row_fill_alt=LIGHT_GREY,
          font_size=11, header_font_size=11):
    """Draw a simple table with coloured header row."""
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [w / n_cols] * n_cols

    row_h = Inches(0.33)
    header_h = Inches(0.38)

    # Header row
    x = l
    for i, (hdr, cw) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, t, cw, header_h, fill=header_fill)
        add_text(slide, hdr, x + Inches(0.06), t + Pt(4), cw - Inches(0.1), header_h,
                 size=header_font_size, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        x += cw

    # Data rows
    for r_idx, row in enumerate(rows):
        fill = row_fill_alt if r_idx % 2 == 0 else WHITE
        y = t + header_h + r_idx * row_h
        x = l
        for i, (cell, cw) in enumerate(zip(row, col_widths)):
            add_rect(slide, x, y, cw, row_h, fill=fill, line=RGBColor(0xDD, 0xDD, 0xDD))
            cell_color = DARK_GREY
            if str(cell).startswith("✅"):
                cell_color = GREEN
            elif str(cell).startswith("⏳"):
                cell_color = AMBER
            add_text(slide, str(cell), x + Inches(0.06), y + Pt(3),
                     cw - Inches(0.1), row_h,
                     size=font_size, color=cell_color, align=PP_ALIGN.LEFT)
            x += cw


def bullet_block(slide, items, l, t, w, h, size=13, indent=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = DARK_GREY
        if indent:
            p.level = 1


def label(slide, text, l, t, w, h, bg=CMU_RED):
    add_rect(slide, l, t, w, h, fill=bg)
    add_text(slide, text, l + Inches(0.08), t + Pt(3), w - Inches(0.1), h,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Total slide count (set once known)
# ---------------------------------------------------------------------------
TOTAL = 11


# ===========================================================================
# SLIDE 1 — Title
# ===========================================================================
slide = prs.slides.add_slide(BLANK)

add_rect(slide, 0, 0, W, H, fill=CMU_RED)
add_rect(slide, Inches(0.5), Inches(2.0), Inches(12.33), Inches(3.5), fill=WHITE)

add_text(slide, "Data Collection Summary",
         Inches(0.8), Inches(2.25), Inches(11.5), Inches(1.2),
         size=40, bold=True, color=CMU_RED, align=PP_ALIGN.LEFT)
add_text(slide, "Stablecoin Depeg Prediction at 5-Minute Resolution",
         Inches(0.8), Inches(3.35), Inches(11.5), Inches(0.6),
         size=20, bold=False, color=DARK_GREY, align=PP_ALIGN.LEFT)
add_text(slide, "CMU MSBA Capstone  ·  2026",
         Inches(0.8), Inches(4.85), Inches(11.5), Inches(0.4),
         size=14, bold=False, color=MID_GREY, align=PP_ALIGN.LEFT)
add_text(slide, "7 sources  ·  7 stablecoins  ·  5-minute resolution",
         Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.4),
         size=13, bold=False, color=MID_GREY, align=PP_ALIGN.LEFT)

footer(slide, 1, TOTAL)


# ===========================================================================
# SLIDE 2 — Stablecoins Covered
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Stablecoins Covered", "7 coins including active, discontinued, and failed examples")

table(slide,
      ["Coin", "Full Name", "Type", "Status", "Coverage Start", "Coverage End"],
      [
          ["USDT",  "Tether",       "Fiat-backed",           "Active",        "2015-01-01", "present"],
          ["USDC",  "USD Coin",     "Fiat-backed",           "Active",        "2018-09-01", "present"],
          ["DAI",   "Dai",          "Crypto-collateralized", "Active",        "2017-12-01", "present"],
          ["BUSD",  "Binance USD",  "Fiat-backed",           "Discontinued",  "2019-09-01", "2023-03-31"],
          ["UST",   "TerraUSD",     "Algorithmic",           "Failed ⚠",     "2020-09-01", "2022-05-31"],
          ["USDe",  "Ethena USDe",  "Synthetic",             "Active",        "2024-02-01", "present"],
          ["RLUSD", "Ripple USD",   "Fiat-backed",           "Active",        "2024-12-01", "present"],
      ],
      l=Inches(0.4), t=Inches(1.3), w=Inches(12.5),
      col_widths=[Inches(1.0), Inches(2.0), Inches(2.5), Inches(1.8), Inches(2.0), Inches(2.0)],
      font_size=12)

add_text(slide,
         "Failed and discontinued coins are intentionally included — they provide real depeg training examples (UST 2022, BUSD 2023).",
         Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.45),
         size=11, italic=True, color=MID_GREY)

footer(slide, 2, TOTAL)


# ===========================================================================
# SLIDE 3 — Architecture Overview
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Architecture Overview", "All sources resolve to 5-minute UTC bins")

# Two columns: native 5m | daily forward-filled
add_rect(slide, Inches(0.4), Inches(1.25), Inches(6.0), Inches(0.38), fill=BLUE)
add_text(slide, "Native 5-Minute Sources", Inches(0.5), Inches(1.28), Inches(5.8), Inches(0.35),
         size=13, bold=True, color=WHITE)

native = [
    "Binance  —  free OHLCV (BTCUSDT, ETHUSDT, stablecoin pairs)",
    "CoinAPI  —  paid VWAP index (primary peg measurement)",
    "On-chain ETH  —  mint/burn + USDT treasury flows (Etherscan V2)",
    "USDT TRON  —  treasury inflow/outflow (TronGrid)",
    "Curve pools  —  TokenExchange swaps (3pool, USDe/USDC, RLUSD/USDC)",
]
bullet_block(slide, native, Inches(0.5), Inches(1.7), Inches(5.7), Inches(2.8), size=12)

add_rect(slide, Inches(6.9), Inches(1.25), Inches(6.0), Inches(0.38), fill=MID_GREY)
add_text(slide, "Daily Sources (forward-filled → 5m)", Inches(7.0), Inches(1.28), Inches(5.8), Inches(0.35),
         size=13, bold=True, color=WHITE)

daily = [
    "FRED  —  DXY, VIX, T10Y, Fed Funds (macro regime)",
    "Market  —  BTC/ETH prices, Fear & Greed index",
]
bullet_block(slide, daily, Inches(7.0), Inches(1.7), Inches(5.7), Inches(1.2), size=12)

# Arrow / merge box
add_rect(slide, Inches(3.5), Inches(4.8), Inches(6.3), Inches(0.5), fill=CMU_RED)
add_text(slide, "merge_sources.py  →  data/processed/{coin}_5m.parquet",
         Inches(3.6), Inches(4.85), Inches(6.1), Inches(0.4),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide,
         "Depeg threshold: ±0.5% from $1.00 peg  ·  Base interval: 5 minutes  ·  Label source: CoinAPI VWAP",
         Inches(0.4), Inches(5.6), Inches(12.5), Inches(0.4),
         size=11, italic=True, color=MID_GREY)

footer(slide, 3, TOTAL)


# ===========================================================================
# SLIDE 4 — Binance + CoinAPI
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Price & OHLCV Data", "Binance (free) and CoinAPI VWAP Index (paid)")

# Left: Binance
label(slide, "1  Binance OHLCV", Inches(0.4), Inches(1.3), Inches(5.8), Inches(0.35))
bullet_block(slide, [
    "Endpoint:  GET https://api.binance.com/api/v3/klines",
    "Params:  symbol, interval=5m, startTime, endTime, limit=1000",
    "Auth:  None (public)",
    "Coverage:  Aug 2017 → present",
], Inches(0.5), Inches(1.75), Inches(5.6), Inches(1.4), size=11)

table(slide,
      ["File", "Rows", "Range"],
      [
          ["usdt_btcusdt", "895,977", "2017-08 → 2026-02"],
          ["usdt_ethusdt", "895,978", "2017-08 → 2026-02"],
          ["usdt_usdcusdt", "709,267", "2018-12 → 2026-02"],
      ],
      l=Inches(0.4), t=Inches(3.3), w=Inches(5.8),
      col_widths=[Inches(2.6), Inches(1.2), Inches(2.0)],
      font_size=11)

add_text(slide, "Use: BTC/ETH as market regime; stablecoin pairs as relative price signal",
         Inches(0.4), Inches(4.55), Inches(5.8), Inches(0.5),
         size=11, italic=True, color=MID_GREY)

# Right: CoinAPI
label(slide, "2  CoinAPI VWAP Index", Inches(6.9), Inches(1.3), Inches(6.0), Inches(0.35), bg=BLUE)
bullet_block(slide, [
    "Endpoint:  GET https://rest-api.indexes.coinapi.io/v1/",
    "             indexes/{index_id}/timeseries",
    "Params:  period_id=5MIN, time_start, time_end, limit=100000",
    "Auth:  X-CoinAPI-Key header  (COINAPI_KEY)",
    "Coverage:  Jun 2017 → present (varies by coin)",
], Inches(7.0), Inches(1.75), Inches(5.8), Inches(1.6), size=11)

table(slide,
      ["Index", "Rows", "Range"],
      [
          ["IDX_REFRATE_VWAP_USDT", "916,454", "2017-06 → 2026-02"],
          ["IDX_REFRATE_VWAP_USDC", "772,132", "2018-10 → 2026-02"],
          ["IDX_REFRATE_VWAP_DAI",  "825,612", "2018-04 → 2026-02"],
          ["IDX_REFRATE_VWAP_BUSD", "527,339", "2019-09 → 2026-02"],
          ["IDX_REFRATE_VWAP_UST",  "528,675", "2020-11 → 2025-12"],
          ["IDX_REFRATE_VWAP_USDE", "198,329", "2024-04 → 2026-02"],
          ["IDX_REFRATE_VWAP_RLUSD","98,904",  "2025-02 → 2026-02"],
      ],
      l=Inches(6.9), t=Inches(3.3), w=Inches(6.0),
      col_widths=[Inches(2.8), Inches(1.2), Inches(2.0)],
      header_fill=BLUE, font_size=11)

footer(slide, 4, TOTAL)


# ===========================================================================
# SLIDE 5 — On-chain ETH Mint/Burn
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "On-Chain Ethereum: Mint / Burn Events",
           "Etherscan V2 getLogs API  ·  Adaptive chunk sizing  ·  100k blocks/chunk")

bullet_block(slide, [
    "Endpoint:  GET https://api.etherscan.io/v2/api?chainid=1&module=logs&action=getLogs",
    "Params:  address (contract), topic0 (event hash), fromBlock, toBlock, offset=1000",
    "Auth:  apikey param  (ETHERSCAN_API_KEY)  ·  Rate: 1 call / 0.25s",
], Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.9), size=11)

table(slide,
      ["Coin", "Mint Event", "Burn Event", "Raw Events", "5m Rows", "Range"],
      [
          ["USDT",  "Issue(uint256)",               "DestroyedBlackFunds",      "11,215",    "8,703",   "2017-11 → 2026-02"],
          ["USDC",  "Mint(addr,addr,uint256)",       "Burn(addr,uint256)",       "2,209,692", "390,104", "2018-09 → 2026-03"],
          ["BUSD",  "SupplyIncreased(addr,uint256)", "SupplyDecreased",          "3,767",     "3,765",   "2019-09 → 2026-01"],
          ["DAI",   "Transfer(from=0x0,...)",         "Transfer(...,to=0x0)",    "1,400,464", "425,580", "2019-11 → 2026-03"],
          ["USDe",  "Transfer(from=0x0,...)",         "Transfer(...,to=0x0)",    "38,001",    "24,589",  "2023-11 → 2026-03"],
          ["RLUSD", "Transfer(from=0x0,...)",         "Transfer(...,to=0x0)",    "288",       "282",     "2024-08 → 2026-02"],
      ],
      l=Inches(0.4), t=Inches(2.2), w=Inches(12.5),
      col_widths=[Inches(1.0), Inches(2.6), Inches(2.3), Inches(1.4), Inches(1.2), Inches(2.5)],
      font_size=10.5)

add_text(slide, "5m bin columns:  mint_count, mint_volume_usd, burn_count, burn_volume_usd, net_flow_usd",
         Inches(0.4), Inches(5.05), Inches(12.5), Inches(0.35), size=11, bold=True, color=DARK_GREY)

add_text(slide,
         "⚠  USDT: Tether never calls Redeem() on Ethereum. Burns are only via DestroyedBlackFunds (OFAC sanctions). "
         "Real redemption pressure tracked separately via treasury wallet flows (next slide).",
         Inches(0.4), Inches(5.5), Inches(12.5), Inches(0.65),
         size=11, italic=True, color=AMBER)

footer(slide, 5, TOTAL)


# ===========================================================================
# SLIDE 6 — USDT Treasury (ETH + TRON)
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "USDT Treasury Flows  —  Leading Indicator",
           "Institutions redeem at par directly with Tether before market depegs")

# Hypothesis box
add_rect(slide, Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.65), fill=LIGHT_GREY, line=CMU_RED)
add_text(slide,
         "Hypothesis: sophisticated institutions send USDT to Tether's treasury wallet "
         "for off-chain USD redemption before the open-market price breaks. "
         "Treasury inflow spikes may precede peg deviations by minutes to hours.",
         Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.6),
         size=11, italic=True, color=DARK_GREY)

# Left: ETH treasury
label(slide, "3  ETH Treasury  (Etherscan V2)", Inches(0.4), Inches(2.05), Inches(6.0), Inches(0.35))
bullet_block(slide, [
    "Wallet:  0x5754284f345afc66a98fbb0a0afe71e0f007b949",
    "Tracks ERC-20 Transfer events TO/FROM the treasury",
    "treasury_inflow  →  institution sends USDT (redemption pressure)",
    "treasury_outflow →  treasury re-issues USDT to institution",
    "Included in usdt_eth_events.parquet / usdt_eth_5m.parquet",
], Inches(0.5), Inches(2.5), Inches(5.7), Inches(2.0), size=11)

# Right: TRON treasury
label(slide, "4  TRON Treasury  (TronGrid)", Inches(6.9), Inches(2.05), Inches(6.0), Inches(0.35), bg=BLUE)
bullet_block(slide, [
    "Endpoint:  GET https://api.trongrid.io/v1/accounts/{addr}/transactions/trc20",
    "Params:  contract_address (USDT TRC20), min/max_timestamp, limit=200, fingerprint",
    "Auth:  TRON-PRO-API-KEY header (optional)",
    "TRON holds >50% of all USDT supply",
    "Inter-treasury transfers excluded (deduplication)",
], Inches(7.0), Inches(2.5), Inches(5.7), Inches(1.85), size=11)

table(slide,
      ["Chain", "Wallet / Scope", "Transfers", "5m Rows", "Range"],
      [
          ["ETH",  "0x5754284f...007b949",          "~11K events",  "8,703",  "2017-11 → 2026-02"],
          ["TRON", "3 treasury wallets (2019–now)", "12,663",       "9,610",  "2019-04 → 2026-02"],
      ],
      l=Inches(0.4), t=Inches(4.6), w=Inches(12.5),
      col_widths=[Inches(1.0), Inches(4.0), Inches(2.0), Inches(1.5), Inches(3.0)],
      font_size=11)

footer(slide, 6, TOTAL)


# ===========================================================================
# SLIDE 7 — Curve Pool Swaps
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Curve Pool Swaps  —  DEX Stress Indicator  ✅ Complete",
           "TokenExchange events  ·  Etherscan V2  ·  655,093 total events across 3 pools")

# Rationale
add_rect(slide, Inches(0.4), Inches(1.25), Inches(12.5), Inches(0.6), fill=LIGHT_GREY, line=CMU_RED)
add_text(slide,
         "Signal: when usdt_net_sell_volume_usd > 0 in the 3pool, traders are selling USDT for DAI/USDC on-chain "
         "— a leading indicator that may precede open-market depeg by minutes. "
         "Observed during UST 2022 collapse and USDC SVB crisis (March 2023).",
         Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.55),
         size=11, italic=True, color=DARK_GREY)

# API
bullet_block(slide, [
    "Endpoint:  GET https://api.etherscan.io/v2/api?chainid=1&module=logs&action=getLogs",
    "topic0:  0x8b3e96f2b889fa771c53c981b40daf005f63f637f1869f707052d15a3dd97140",
    "         keccak256(\"TokenExchange(address,int128,uint256,int128,uint256)\")  —  same for all Curve pool types",
    "Data field:  [32B: sold_id][32B: tokens_sold][32B: bought_id][32B: tokens_bought]",
], Inches(0.4), Inches(2.0), Inches(12.5), Inches(1.1), size=11)

table(slide,
      ["Pool", "Contract Address", "Tokens", "Deployed", "Status"],
      [
          ["3pool",      "0xbEbc44782C7dB0a1A60Cb6fe97d0b483032FF1C7", "DAI / USDC / USDT", "Sep 2020", "✅ 538,925 events / 277,828 5m rows"],
          ["usde_usdc",  "0x02950460E2b9529D0E00284A5fA2D7Bdf3Fa4d72", "USDe / USDC",       "Nov 2023", "✅ 106,449 events / 61,095 5m rows"],
          ["rlusd_usdc", "0xd001ae433f254283fece51d4acce8c53263aa186", "RLUSD / USDC",      "Dec 2024", "✅ 10,719 events / 8,835 5m rows"],
      ],
      l=Inches(0.4), t=Inches(3.25), w=Inches(12.5),
      col_widths=[Inches(1.5), Inches(4.2), Inches(2.2), Inches(1.4), Inches(2.2)],
      font_size=11)

add_text(slide,
         "5m bin columns per token:  {token}_sold_count,  {token}_sold_volume_usd,  "
         "{token}_bought_count,  {token}_bought_volume_usd,  {token}_net_sell_volume_usd",
         Inches(0.4), Inches(4.6), Inches(12.5), Inches(0.4),
         size=11, bold=True, color=DARK_GREY)

footer(slide, 7, TOTAL)


# ===========================================================================
# SLIDE 8 — FRED + Market
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Macro & Market Context", "Daily sources forward-filled to 5-minute bins in merge step")

# Left: FRED
label(slide, "6  FRED Macro", Inches(0.4), Inches(1.3), Inches(5.8), Inches(0.35))
bullet_block(slide, [
    "Endpoint:  GET https://api.stlouisfed.org/fred/series/observations",
    "Params:  series_id, observation_start, file_type=json",
    "Auth:  api_key param  (FRED_API_KEY)",
    "Frequency:  daily → forward-filled to 5m",
    "Coverage:  2015-01-01 → 2026-02-19  (2,885 rows)",
], Inches(0.5), Inches(1.75), Inches(5.6), Inches(1.7), size=11)

table(slide,
      ["Series ID", "Description"],
      [
          ["DTWEXBGS", "USD Index (DXY) — dollar strength"],
          ["VIXCLS",   "CBOE VIX — equity volatility / fear"],
          ["DGS10",    "10-yr Treasury yield"],
          ["FEDFUNDS", "Fed Funds rate"],
      ],
      l=Inches(0.4), t=Inches(3.6), w=Inches(5.8),
      col_widths=[Inches(1.8), Inches(4.0)],
      font_size=11)

# Right: Market
label(slide, "7  Market Daily  (CoinGecko + Alternative.me)", Inches(6.9), Inches(1.3), Inches(6.0), Inches(0.35), bg=BLUE)
bullet_block(slide, [
    "CoinGecko:  GET https://api.coingecko.com/api/v3/coins/{id}/market_chart",
    "   Params:  vs_currency=usd, days, interval=daily",
    "   Auth:  x-cg-demo-api-key header (COINGECKO_API_KEY, optional)",
    "Alternative.me:  GET https://api.alternative.me/fng/",
    "   Params:  limit, format=json  (no auth required)",
    "Frequency:  daily → forward-filled to 5m",
    "Coverage:  2018-02-01 → 2026-02-22  (2,940 rows)",
], Inches(7.0), Inches(1.75), Inches(5.7), Inches(2.4), size=11)

table(slide,
      ["Feature", "Source", "Description"],
      [
          ["btc_price",    "CoinGecko", "Bitcoin USD price"],
          ["eth_price",    "CoinGecko", "Ethereum USD price"],
          ["fear_greed",   "Alternative.me", "0–100 sentiment index"],
      ],
      l=Inches(6.9), t=Inches(4.3), w=Inches(6.0),
      col_widths=[Inches(1.5), Inches(2.0), Inches(2.5)],
      header_fill=BLUE, font_size=11)

footer(slide, 8, TOTAL)


# ===========================================================================
# SLIDE 9 — Collection Status
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Collection Status", f"As of 2026-03-03  —  All sources complete")

table(slide,
      ["Source", "Script", "Status", "Coverage", "Notes"],
      [
          ["Binance OHLCV",        "collect_binance.py",   "✅ Complete", "Through 2026-02-28", "Free, no key"],
          ["CoinAPI VWAP",         "collect_coinapi.py",   "✅ Complete", "Through 2026-02-21", "Paid key required"],
          ["ETH mint/burn",        "collect_onchain.py",   "✅ Complete", "Through 2026-03-02", "6 coins, 3.7M events"],
          ["USDT ETH treasury",    "collect_onchain.py",   "✅ Complete", "Through 2026-02-28", "Incl. in USDT onchain"],
          ["USDT TRON treasury",   "collect_tron.py",      "✅ Complete", "Through 2026-02-28", "12,663 transfers"],
          ["Curve pool swaps",     "collect_curve.py",     "✅ Complete", "Through 2026-03-03", "655K events, 3 pools"],
          ["FRED macro",           "collect_fred.py",      "✅ Complete", "Through 2026-02-19", "4 series"],
          ["Market daily",         "collect_market.py",    "✅ Complete", "Through 2026-02-22", "BTC/ETH/Fear&Greed"],
      ],
      l=Inches(0.4), t=Inches(1.3), w=Inches(12.5),
      col_widths=[Inches(2.2), Inches(2.3), Inches(1.7), Inches(2.4), Inches(3.3)],
      font_size=11)

add_rect(slide, Inches(0.4), Inches(5.85), Inches(12.5), Inches(0.55), fill=LIGHT_GREY, line=CMU_RED)
add_text(slide,
         "Next step:  Once Curve collection completes, run  merge_sources.py  "
         "to build  data/processed/{coin}_5m.parquet  per-coin feature sets.",
         Inches(0.55), Inches(5.9), Inches(12.2), Inches(0.5),
         size=12, bold=True, color=DARK_GREY)

footer(slide, 9, TOTAL)


# ===========================================================================
# SLIDE 10 — Feature Taxonomy
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "Feature Taxonomy", "Signals extracted from each source after merge")

table(slide,
      ["Category", "Features", "Source(s)", "Frequency"],
      [
          ["Peg price",           "VWAP vs $1.00, bid/ask spread, OHLC",        "CoinAPI, Binance",         "5m (native)"],
          ["On-chain supply",     "mint_volume_usd, burn_volume_usd, net_flow",  "ETH on-chain",             "5m (aggregated)"],
          ["Redemption pressure", "treasury_inflow/outflow (ETH + TRON)",        "ETH on-chain, TronGrid",   "5m (aggregated)"],
          ["DEX stress",          "{token}_net_sell_volume_usd per Curve pool",  "Curve (Etherscan V2)",     "5m (aggregated)"],
          ["Market regime",       "btc_return, eth_return, fear_greed_index",    "Binance, CoinGecko",       "5m (fwd-filled)"],
          ["Macro",               "DXY, VIX, T10Y_yield, fed_funds_rate",        "FRED",                     "5m (fwd-filled)"],
      ],
      l=Inches(0.4), t=Inches(1.3), w=Inches(12.5),
      col_widths=[Inches(2.4), Inches(4.2), Inches(3.0), Inches(2.3)],
      font_size=12)

add_text(slide, "Depeg label:  VWAP price deviation > ±0.5% from $1.00  →  binary classification target",
         Inches(0.4), Inches(4.8), Inches(12.5), Inches(0.4),
         size=13, bold=True, color=CMU_RED)

add_text(slide,
         "Leading indicator signals (Curve, treasury flows) are the key differentiator from purely price-based models. "
         "These capture informed on-chain behaviour before CEX prices move.",
         Inches(0.4), Inches(5.3), Inches(12.5), Inches(0.55),
         size=11, italic=True, color=MID_GREY)

footer(slide, 10, TOTAL)


# ===========================================================================
# SLIDE 11 — API Keys Reference
# ===========================================================================
slide = prs.slides.add_slide(BLANK)
header_bar(slide, "API Keys & Rate Limits", "Set in .env file at project root")

table(slide,
      ["Env Var", "Service", "Required?", "Rate Limit", "Used For"],
      [
          ["ETHERSCAN_API_KEY",  "etherscan.io",          "Required",  "5 calls/sec",        "ETH on-chain events, Curve"],
          ["COINAPI_KEY",        "coinapi.io (Indexes)",  "Required",  "Varies by plan",     "VWAP index data (paid)"],
          ["FRED_API_KEY",       "fred.stlouisfed.org",   "Required",  "120 calls/min",      "Macro indicators"],
          ["TRONGRID_API_KEY",   "trongrid.io",           "Optional",  "20 QPS with key",    "TRON treasury (10 QPS without)"],
          ["COINGECKO_API_KEY",  "coingecko.com",         "Optional",  "30 calls/min (free)","Market daily (higher limit)"],
      ],
      l=Inches(0.4), t=Inches(1.3), w=Inches(12.5),
      col_widths=[Inches(2.5), Inches(2.3), Inches(1.4), Inches(2.0), Inches(4.3)],
      font_size=12)

add_text(slide, "Pagination strategies:",
         Inches(0.4), Inches(4.2), Inches(12.5), Inches(0.35),
         size=12, bold=True, color=DARK_GREY)

bullet_block(slide, [
    "Binance / CoinAPI:  cursor-based on time (startTime / time_start) — 1,000 / 100,000 results per page",
    "Etherscan:  block range chunking (100k blocks/chunk) with adaptive bisection when 1,000-result limit hit",
    "TronGrid:  fingerprint token in meta field for next-page cursor within each 30-day time window",
    "FRED:  single request per series (full history returned at once)",
], Inches(0.5), Inches(4.6), Inches(12.2), Inches(1.8), size=11)

footer(slide, 11, TOTAL)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out = Path("docs/data_collection_summary.pptx")
prs.save(out)
print(f"Saved {out}  ({TOTAL} slides)")
